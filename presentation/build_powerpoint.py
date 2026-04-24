"""Build the PRISM final PowerPoint deck from local release artifacts."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "PRISM_final_presentation.pptx"
ASSETS = ROOT / "presentation" / "assets"

INK = RGBColor(23, 23, 23)
MUTED = RGBColor(88, 91, 99)
ACCENT = RGBColor(11, 111, 106)
ACCENT_2 = RGBColor(155, 61, 18)
ACCENT_3 = RGBColor(52, 89, 149)
BG = RGBColor(247, 245, 239)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(217, 212, 199)
SOFT = RGBColor(233, 242, 239)
WARN = RGBColor(255, 243, 216)


def set_bg(slide, color=BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, x, y, w, h, size=22, bold=False, color=INK, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def kicker(slide, text, y=0.42):
    return textbox(slide, text.upper(), 0.62, y, 11.6, 0.28, size=11, bold=True, color=ACCENT)


def title(slide, text, y=0.78, size=42):
    return textbox(slide, text, 0.62, y, 11.8, 0.9, size=size, bold=True, color=INK)


def body(slide, text, x=0.65, y=1.75, w=8.7, h=1.0, size=19):
    return textbox(slide, text, x, y, w, h, size=size, color=MUTED)


def card(slide, x, y, w, h, heading, text, accent=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    textbox(slide, heading, x + 0.18, y + 0.18, w - 0.36, 0.34, size=16, bold=True, color=accent)
    textbox(slide, text, x + 0.18, y + 0.62, w - 0.36, h - 0.78, size=12.5, color=MUTED)


def metric_card(slide, x, y, w, h, metric, label, accent=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    textbox(slide, metric, x + 0.18, y + 0.14, w - 0.36, 0.45, size=25, bold=True, color=accent)
    textbox(slide, label, x + 0.18, y + 0.68, w - 0.36, h - 0.76, size=11.5, color=MUTED)


def pill(slide, text, x, y, w, color=SOFT, text_color=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = LINE
    tb = textbox(slide, text, x + 0.08, y + 0.06, w - 0.16, 0.24, size=10.5, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return tb


def add_image(slide, path, x, y, w, h=None):
    kwargs = {"left": Inches(x), "top": Inches(y), "width": Inches(w)}
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), **kwargs)


def footer(slide, idx):
    textbox(slide, "PRISM | CSE 579 Final Presentation", 0.62, 7.05, 4.4, 0.22, size=9, color=MUTED)
    textbox(slide, f"{idx} / 7", 12.0, 7.05, 0.7, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=17):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(6)
    return shape


def small_contribution_bar(slide, y=6.56):
    textbox(
        slide,
        "Contributions: Ritik - production integration, demo, report, release package | Omkar - TBD | Vaishnavi - TBD | Moin - TBD",
        0.62,
        y,
        12.0,
        0.28,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "CSE 579 Knowledge Representation and Reasoning")
    textbox(slide, "PRISM", 0.62, 1.05, 8.6, 0.9, size=58, bold=True)
    body(slide, "Representation-aware retrieval routing for structural mismatch in question answering.", y=2.08, w=10.4, size=23)
    pill(slide, "Production: computed_ras", 0.65, 3.05, 2.45, SOFT, ACCENT)
    pill(slide, "Research overlays: rescue, RAS_V3, RAS_V4", 3.25, 3.05, 3.85, WARN, ACCENT_2)
    pill(slide, "BM25 + Dense + KG + Hybrid", 7.25, 3.05, 2.8, RGBColor(235, 239, 248), ACCENT_3)
    card(slide, 0.65, 4.05, 3.2, 1.25, "Problem", "Retrieval can fail when the representation does not match the query.")
    card(slide, 4.05, 4.05, 3.2, 1.25, "KRR approach", "Reason over query structure, route adequacy, evidence, and trace.")
    card(slide, 7.45, 4.05, 3.2, 1.25, "Contribution", "An inspectable router across BM25, Dense, KG, and Hybrid.")
    small_contribution_bar(slide)
    footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "Problem")
    title(slide, "One retriever cannot preserve every kind of structure.", size=36)
    body(slide, "A query can be answerable in the corpus but fail because the retrieval representation is inadequate.", y=1.62, w=10.6, size=19)
    card(slide, 0.65, 2.55, 2.75, 1.55, "Lexical", "RFC numbers, identifiers, and code-like strings need exact matching.", ACCENT_2)
    card(slide, 3.7, 2.55, 2.75, 1.55, "Semantic", "Paraphrases and conceptual similarity need dense retrieval.", ACCENT_3)
    card(slide, 6.75, 2.55, 2.75, 1.55, "Deductive", "Class, property, and membership questions need structured graph evidence.", ACCENT)
    card(slide, 9.8, 2.55, 2.75, 1.55, "Relational", "Bridge questions need fused text and structural evidence.", RGBColor(95, 78, 160))
    textbox(slide, "PRISM makes representation choice explicit, inspectable, and measurable.", 0.65, 5.0, 11.4, 0.5, size=24, bold=True, color=INK)
    card(slide, 0.65, 5.85, 11.4, 0.62, "Motivating example", "The same corpus can contain the answer, but RFC-7231, climate-anxiety paraphrases, and mammal capability questions require different representations.", ACCENT_3)
    footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "KRR insight")
    title(slide, "Structural mismatch is a representation adequacy problem.", size=35)
    body(slide, "The reasoning step is not only answer generation. PRISM first reasons about which representation can support the query.", y=1.55, w=11.3, size=18)
    card(slide, 0.75, 2.55, 2.55, 1.45, "1. Query", "Extract lexical, semantic, deductive, and relational signals.")
    card(slide, 3.55, 2.55, 2.55, 1.45, "2. Route", "RAS scores BM25, Dense, KG, and Hybrid candidates.", ACCENT_3)
    card(slide, 6.35, 2.55, 2.55, 1.45, "3. Evidence", "Selected backend retrieves provenance-bearing support.", ACCENT_2)
    card(slide, 9.15, 2.55, 2.55, 1.45, "4. Trace", "Answer cites evidence ids and exposes rationale.", ACCENT)
    add_image(slide, ASSETS / "ras_family_overview.png", 1.65, 4.6, 9.9, 1.65)
    textbox(slide, "This is the reasoning process behind the method: represent the query need, choose the adequate representation, then justify the answer with evidence.", 1.25, 6.28, 10.6, 0.34, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "System")
    title(slide, "PRISM architecture and RAS", size=38)
    add_image(slide, ASSETS / "architecture_diagram.png", 0.65, 1.62, 6.15, 4.6)
    card(slide, 7.1, 1.75, 4.9, 0.9, "Production", "computed_ras is the default deterministic router.", ACCENT)
    card(slide, 7.1, 2.9, 4.9, 0.9, "RAS_V3", "Interpretable route-adequacy model; analysis-only.", ACCENT_3)
    card(slide, 7.1, 4.05, 4.9, 0.9, "RAS_V4", "Joint route-and-evidence adequacy model; analysis-only.", ACCENT_3)
    card(slide, 7.1, 5.2, 4.9, 0.9, "Calibrated rescue", "Research/demo overlay strongest on adversarial answer accuracy.", ACCENT_2)
    textbox(slide, "Design choice: keep production deterministic and use research overlays to measure hard-case improvements without changing the default route policy.", 7.15, 6.28, 4.8, 0.34, size=11.5, color=MUTED)
    footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "Evaluation")
    title(slide, "Strong on stable layers, honest on hard cases.", size=34)
    add_image(slide, ASSETS / "benchmark_overview.png", 0.65, 1.55, 5.8, 3.25)
    add_image(slide, ASSETS / "adversarial_router_comparison.png", 6.75, 1.55, 5.8, 3.25)
    metric_card(slide, 0.8, 5.35, 2.8, 1.0, "80/80", "Curated end-to-end")
    metric_card(slide, 3.9, 5.35, 2.8, 1.0, "1.000", "Public raw and public graph tests")
    metric_card(slide, 7.0, 5.35, 2.8, 1.0, "0.917", "computed_ras adversarial test answer")
    metric_card(slide, 10.1, 5.35, 2.8, 1.0, "0.958", "calibrated rescue adversarial test answer", ACCENT_2)
    textbox(slide, "Finding: representation-aware routing is stable on standard layers; hard ambiguity shows evidence rescue is complementary.", 1.2, 6.55, 10.8, 0.28, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "Demo and evidence")
    title(slide, "The UI exposes the full reasoning path.", size=36)
    bullet_list(
        slide,
        [
            "Demo path: Executive Summary -> Guided Demo -> Demo / Query.",
            "Canonical queries cover lexical, semantic, deductive, and relational retrieval.",
            "Each run shows parsed features, route scores, selected backend, evidence ids, answer, and trace.",
            "Human evaluation supports trace clarity and faithfulness, with small-study caveats.",
        ],
        0.85,
        1.8,
        5.45,
        3.6,
        size=17,
    )
    add_image(slide, ASSETS / "human_eval_overview.png", 6.8, 1.65, 5.4, 3.5)
    card(slide, 6.95, 5.45, 5.1, 0.75, "Demo fallback", "Benchmark-safe mode remains reproducible if optional open-corpus or LLM components are unavailable.", ACCENT_2)
    footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    kicker(slide, "Takeaways")
    title(slide, "Takeaways, contributions, and limitations", size=35)
    card(slide, 0.75, 1.55, 3.55, 1.32, "Main contribution", "Representation-aware routing across BM25, Dense, KG, and Hybrid.")
    card(slide, 4.85, 1.55, 3.55, 1.32, "Main result", "Stable benchmark performance; calibrated rescue is strongest on hard adversarial answer accuracy.", ACCENT_3)
    card(slide, 8.95, 1.55, 3.55, 1.32, "Main limitation", "Hard route-boundary cases still require better uncertainty and evidence adequacy.", ACCENT_2)
    card(slide, 0.75, 3.35, 2.85, 1.1, "Ritik", "Production integration, routing/demo stability, final report, release package, presentation deliverables.", ACCENT)
    card(slide, 3.9, 3.35, 2.85, 1.1, "Omkar", "TBD by teammate.", ACCENT_3)
    card(slide, 7.05, 3.35, 2.85, 1.1, "Vaishnavi", "TBD by teammate.", ACCENT)
    card(slide, 10.2, 3.35, 2.85, 1.1, "Moin", "TBD by teammate.", ACCENT_2)
    textbox(slide, "Final decision: production remains computed_ras. PRISM is bounded source-pack/local/open-corpus QA, not web-scale search.", 0.85, 5.35, 11.35, 0.7, size=21, bold=True, color=INK, align=PP_ALIGN.CENTER)
    textbox(slide, "Presentation target: finish in 5 minutes; each slide maps directly to problem, KRR approach, reasoning process, results, and contribution.", 1.1, 6.38, 10.9, 0.28, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 7)

    prs.save(OUT)


if __name__ == "__main__":
    build()
    print(f"wrote {OUT}")
